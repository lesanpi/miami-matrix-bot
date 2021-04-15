from __future__ import print_function, unicode_literals
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
# from selenium.webdriver.common.action_chains import ActionChains
from openpyxl import load_workbook
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.common.keys import Keys
from webdriver_manager.chrome import ChromeDriverManager
from PIL import Image
from time import sleep
import os
import pprint
from pptx import Presentation
from pptx.util import Inches, Pt
from pyfiglet import Figlet
from PyInquirer import prompt, print_json, style_from_dict, Token, Separator
from clint.textui import colored, puts
from os import system, name
import glob
import traceback

# Limpiar la pantalla
def clear():
    # Windows
    if name == 'nt':
        _ = system('cls')
    # Mac and Linux
    else:
        _ = system('clear')

# Styles
style = style_from_dict({
    Token.Separator: '#cc5454',
    Token.QuestionMark: '#673ab7 bold',
    Token.Selected: '#cc5454',  # default
    Token.Pointer: '#673ab7 bold',
    Token.Instruction: '',  # default
    Token.Answer: '#f44336 bold',
    Token.Question: 'bold #673ab7',
})

# Titulo
title = Figlet(font='slant')
# Actions
main_menu_actions = [
    "Use a especific address",
    "Use a excel file",
]
excel_files = glob.glob('*.xlsx')
# Menu
main_menu = [
    {
        "type": "list",
        "message": "Menu",
        "name": "action",
        "choices": main_menu_actions + [Separator(), "Cerrar"]
    }
]
# Address
address_info_input = [
    {
        "type": "input",
        "message": "Address",
        "name": "address"
    },
    {
        "type": "input",
        "message": "Baths (Press Enter if None)",
        "name": "baths"
    },
    {
        "type": "input",
        "message": "Rooms (Press Enter if None)",
        "name": "rooms"
    },
    {
        "type": "input",
        "message": "Sqft To (Press Enter if None)",
        "name": "sqft_to"
    },
]
# Excel
excel_info_input = [
    {
        "type": "list",
        "message": "Select a Excel File",
        "name": "filename",
        "choices": excel_files
    }
]

printer = pprint.PrettyPrinter(indent=1)

URL = "https://sef.clareityiam.net/idp/login"
COUNTY_URL = "https://www.miamidade.gov/Apps/PA/propertysearch/#/"
GOOGLE_URL = "https://www.google.com"

# Excel Cols
SQFT_COL = 'K'
BEDS_COL = 'L'
BATHS_COL = 'M'
STREET_COL = 'B'
CITY_COL = 'C'
ZIP_COL = 'D'

# Driver
options = webdriver.ChromeOptions()
options.add_argument('--ignore-certificate-errors')
options.add_argument('--ignore-ssl-errors')
#options.add_argument('--headless')
options.add_argument('--log-level=3')
#options.add_argument("--incognito")
driver = None

IDS = {
    "address_input_1": "Fm11_Ctrl7_TB",
    "address_input_2": "Fm16_Ctrl7_TB",
    "address_input_3": "Fm15_Ctrl7_TB",
    "bedrooms_input_1": "Fm11_Ctrl9_TB",
    "bedrooms_input_2": "Fm16_Ctrl9_TB",
    "bedrooms_input_3": "Fm15_Ctrl9_TB",
    "baths_input_1": "Fm11_Ctrl13_TB",
    "baths_input_3": "Fm15_Ctrl13_TB",
    "sqft_living_area_input_1": "Fm11_Ctrl57_TB",
    "sqft_living_area_input_2": "Fm16_Ctrl57_TB",
    "sqft_living_area_input_3": "Fm15_Ctrl57_TB",
    "select_within_1": "Fm11_Ctrl7_Radius",
    "select_within_2": "Fm16_Ctrl7_Radius",
    "select_within_3": "Fm15_Ctrl7_Radius",
    "results_tab": "m_ucResultsPageTabs_m_pnlResultsTab",
    "display": "m_ucDisplayPicker_m_ddlDisplayFormats",
    "read_later": "NewsDetailPostpone",
    "rented_input_3": "FmFm15_Ctrl596_21510_Ctrl596_TB"
}

within_type = 1

XPATHS = {
    "username_input": '//div[@id="clareity"]',
    "password_input": '//div[@id="security"]',
    "login_button": '//button[@id="loginbtn"]',
    "end_tour_button": '//button[@data-role="end"]',
    "matrix_app": '//div[@id="appColumn115"]',
    "search_option": '//a[@href="/Matrix/Search"]/..',
    "options_for_search": '//table[@class="min"]//table[@class="bottom"]//a',
    "dialog_address_search": '//div[@class="mapSearchDialog"]',
    "select_within_option_0.5": f'//select[@id="{IDS[f"select_within_{within_type}"]}"]/option[@value="0.80467200"]',
    "select_within_option_1": f'//select[@id="{IDS[f"select_within_{within_type}"]}"]/option[@value="1.60934400"]',
    "display_closed_comp": f"//select[@id='{IDS['display']}']/option[@value='U73530']",
    "display_marketing_to_realtors": f"//select[@id='{IDS['display']}']/option[@value='U71952']",
    "display_for_sale": f"//select[@id='{IDS['display']}']/option[@value='U74267']",
    "display_ac/p/a_review": f"//select[@id='{IDS['display']}']/option[@value='U84429']",
    "results_table": "//div[@class='css_container']",
    "sp_tab": r"//th[contains(@data-mlheader, '1\bSP$\a2\b')]",
    "current_price_tab": r"//th[@data-mlheader='1\bCurrent Price\a2\bCurrent Price']",
    "distance_tab": r"//th[@data-mlheader='1\bDistance\a2\bDistance']",
    "active_checkbox": '//input[@type="checkbox" and @class="checkbox" and @value="101"]',
    "rented_checkbox": '//input[@type="checkbox" and @class="checkbox" and @value="21510"]',
    "expired_checkbox": '//input[@type="checkbox" and @class="checkbox" and @value="106"]',
    "pending_checkbox": '//input[@type="checkbox" and @class="checkbox" and @value="21508"]',
    "active_with_contract_checkbox": '//input[@type="checkbox" and @class="checkbox" and @value="21505"]',
    "filter_container": "//div[@class='css_container']",
    "folio_number_link": "//div[@class='results_record ng-scope']/div/span"
}


def screenshot_and_crop(folder, location, size, filename):
    verify_folder_exists(folder)
    file_path = f"{folder}/{filename.lower().replace(' ', '-').replace(',', '')}.png"
    driver.save_screenshot(file_path)

    # crop image
    x = location['x']
    y = location['y']
    width = location['x'] + size['width']
    height = location['y'] + size['height']
    im = Image.open(file_path)
    im = im.crop((int(x), int(y), int(width), int(height)))
    im.save(file_path)
    return file_path


def criteria_screenshot(address, folder):
    # Screenshot
    filter_container = WebDriverWait(driver, 15) \
        .until(EC.element_to_be_clickable((By.XPATH, XPATHS['filter_container'])))

    filter_path = screenshot_of_element(filter_container, 'criteria/' + folder, address)
    driver.execute_script("window.scrollTo(0, 0)")
    return filter_path


def screenshot_of_element(element, folder, filename, width=None, height=None):
    location = element.location_once_scrolled_into_view
    size = element.size
    if width:
        size['width'] = width
    if height:
        size['height'] = height
    file_path = screenshot_and_crop(folder, location, size, filename)
    return file_path


def verify_folder_exists(folder):
    if not os.path.exists(folder):
        os.makedirs(folder)


# Login
def login():
    # Login
    driver.get(URL)
    sleep(3)

    for i in range(2):
        # Password
        password_input = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['password_input'])))
        password_input.click()
        password_input.send_keys('1234armando')
        sleep(1)

    sleep(1)
    # Username
    username_input = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['username_input'])))
    username_input.click()
    username_input.send_keys('3457717')

    # Login
    login_button = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['login_button'])))
    login_button.click()


# Select app
def select_matrix_app():
    # Quit dialog
    try:
        WebDriverWait(driver, 15).until(
            EC.element_to_be_clickable((By.XPATH, XPATHS['end_tour_button']))).click()
    except Exception as e:
        print("WARNING: END TOUR Button no necessary.")
    # Select matrix app
    WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['matrix_app']))).click()

    driver.close()
    new_window = driver.window_handles[0]
    driver.switch_to.window(new_window)
    # driver.switch_to_window(new_window)

    try:
        WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.ID, IDS['read_later']))).click()
    except Exception as e:
        print("WARNING: Modal Message Not Found")
        pass


def escape():
    webdriver.ActionChains(driver).send_keys(Keys.ESCAPE).perform()


# Searchs
# Select type of search
def select_search(search: int):
    # Click on search
    search_option = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['search_option'])))
    search_option.click()
    escape()

    # Select search type
    WebDriverWait(driver, 15).until(
        EC.element_to_be_clickable((By.XPATH, XPATHS['options_for_search'])))
    options_for_search = driver.find_elements_by_xpath(XPATHS['options_for_search'])
    options_for_search[search].click()


def single_family_search(address, baths, rooms, sqft_to):
    # Search
    select_search(search=0)

    escape()
    # Filter
    criteria_path = single_family_filter(address, baths, rooms, sqft_to, folder="single_family")
    # Results
    results_path = results_family_search(address)

    return {
        "criteria": criteria_path,
        "results": results_path
    }


def res_income_search(address, baths, rooms, sqft_to):
    # Search
    select_search(search=1)
    escape()
    # Filter
    criteria_path = single_family_filter(address, baths, rooms, sqft_to, miles=1, search_type=2, folder="res_income")

    results_path = results_res_income(address)

    return {
        "criteria": criteria_path,
        "results": results_path
    }


def res_rental_search(address, baths, rooms, sqft_to):
    # Search
    select_search(search=3)
    escape()
    # Filter
    criteria_path = single_family_filter(address, baths, rooms, sqft_to, miles=1, search_type=3, folder="res_rental")
    results_path_1 = results_res_rental(address, display_mode="for_sale")
    results_path_2 = results_res_rental(address, display_mode="marketing")

    return {
        "criteria": criteria_path,
        "results": [results_path_1, results_path_2]
    }


def single_family_search_2(address, baths, rooms, sqft_to):
    # Search
    select_search(search=0)

    escape()
    # Filter
    criteria_path = single_family_filter(address, baths, rooms, sqft_to, search_type=4, folder="single_family_2")
    # Results
    results_path = results_family_search(address, index=2, display_mode='display_ac/p/a_review')

    return {
        "criteria": criteria_path,
        "results": results_path
    }


def res_income_search_2(address, baths, rooms, sqft_to):
    # Search
    select_search(search=1)
    escape()
    # Filter
    criteria_path = single_family_filter(address, baths, rooms, sqft_to, miles=1, search_type=5, folder="res_income_2")
    results_path = results_res_income(address, index=2, display_mode='display_ac/p/a_review')

    return {
        "criteria": criteria_path,
        "results": results_path
    }


# Filters
def single_family_filter(address, baths, rooms, sqft_to, miles=0.5, search_type=1, folder="single_family"):
    # RE1/RE2 Single Family/Condo Filter

    # Change within options
    XPATHS[
        "select_within_option_0.5"] = f'//select[@id="{IDS[f"select_within_{search_type - 3 if search_type > 3 else search_type}"]}"]' \
                                      f'/option[@value="0.80467200"]'
    XPATHS[
        "select_within_option_1"] = f'//select[@id="{IDS[f"select_within_{search_type - 3 if search_type > 3 else search_type}"]}"]/option[@value="1.60934400"]'

    escape()
    # Within open
    select_within = WebDriverWait(driver, 15).until(
        EC.element_to_be_clickable(
            (By.ID, IDS[f'select_within_{search_type - 3 if search_type > 3 else search_type}'])
        )
    )
    select_within.click()

    escape()
    # Select option
    select_within_option = WebDriverWait(driver, 15).until(
        EC.element_to_be_clickable((By.XPATH, XPATHS[f"select_within_option_{miles}"])))

    # Within close
    select_within_option.click()

    # Address input
    address_input = WebDriverWait(driver, 15).until(
        EC.element_to_be_clickable(
            (By.ID, IDS[f'address_input_{search_type - 3 if search_type > 3 else search_type}'])
        )
    )
    escape()
    address_input.click()
    address_input.send_keys(address)

     # Select option
    dialog_address_search = WebDriverWait(driver, 15).until(
        EC.element_to_be_clickable(
            (By.XPATH, XPATHS['dialog_address_search'])
        )
    )
    escape()
    dialog_address_search.click()

    if search_type not in [2, 5]:
        # Bedrooms
        bedrooms_input = WebDriverWait(driver, 15).until(
            EC.element_to_be_clickable(
                (By.ID, IDS[f'bedrooms_input_{search_type - 3 if search_type > 3 else search_type}'])
            )
        )
        if rooms:
            bedrooms_input.send_keys(f"{rooms - 1}-{rooms + 1}")

        # Baths
        baths_input = WebDriverWait(driver, 15).until(
            EC.element_to_be_clickable(
                (By.ID, IDS[f'baths_input_{search_type - 3 if search_type > 3 else search_type}'])
            )
        )
        if baths:
            baths_input.send_keys(f"{baths - 1}-{baths + 1}")

    # SQFT
    sqft_living_area_input = WebDriverWait(driver, 15).until(
        EC.element_to_be_clickable(
            (By.ID, IDS[f'sqft_living_area_input_{search_type - 3 if search_type > 3 else search_type}'])
        )
    )
    if sqft_to:
        sqft_living_area_input.send_keys(f"0-{sqft_to + 700}")

    # Res rental
    if search_type == 3:
        # Active uncheck
        WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['active_checkbox']))).click()
        # Rented check
        WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['rented_checkbox']))).click()
        # Expired check
        WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['expired_checkbox']))).click()

        rented_input = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.ID, IDS['rented_input_3'])))
        rented_input.clear()
        rented_input.send_keys(f"0-365")
    elif search_type == 4 or search_type == 5:
        # Active
        WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['active_checkbox']))).click()
        # Pending
        WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['pending_checkbox']))).click()
        # Active with contract
        WebDriverWait(driver, 15) \
            .until(EC.element_to_be_clickable((By.XPATH, XPATHS['active_with_contract_checkbox']))).click()

    criteria_path = criteria_screenshot(address, folder=folder)

    # Go to Results
    results_tab = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.ID, IDS['results_tab'])))
    results_tab.click()

    return criteria_path


# Results actions
# Order
def order_by_sp():
    sp_tab = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['sp_tab'])))
    sp_tab.click()
    sleep(0.5)

def order_by_current_price():
    current_price_tab = WebDriverWait(driver, 15).until(
        EC.element_to_be_clickable((By.XPATH, XPATHS['current_price_tab'])))
    current_price_tab.click()
    sleep(0.5)

def order_by_distance():
    distance_tab = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['distance_tab'])))
    distance_tab.click()
    sleep(0.5)

# Table actions
def set_display(display_mode='display_closed_comp'):
    # Open dropdown
    display = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.ID, IDS['display'])))
    display.click()
    # Select option
    display_option = WebDriverWait(driver, 15) \
        .until(EC.element_to_be_clickable((By.XPATH, XPATHS[display_mode])))
    display_option.click()
    # Click again to close dropdown
    sleep(2)
    # display.click()


# Results
def results_family_search(address, index=None, display_mode='display_closed_comp'):
    escape()

    set_display(display_mode)
    sleep(3)
    escape()
    try:
        order_by_sp()
        order_by_sp()
    except:
        print("WARNING: Can't Order by SP. Maybe because there is not data")
    sleep(1)
    escape()

    # Screenshot
    results_table = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['results_table'])))

    if not index:
        results_path = screenshot_of_element(results_table, 'results/single_family', address, height=1000)
    else:
        results_path = screenshot_of_element(results_table, f'results/single_family_{index}', address, height=1000)

    driver.execute_script("window.scrollTo(0, 0)")
    return results_path


def results_res_income(address, index=None, display_mode='display_closed_comp'):
    escape()

    set_display(display_mode)
    sleep(1)

    escape()
    try:
        order_by_sp()
        order_by_sp()
    except:
        print("WARNING: Can't Order by SP. Maybe because there is not data")
    sleep(1)
    escape()

    # Screenshot
    results_table = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['results_table'])))

    if not index:
        results_path = screenshot_of_element(results_table, 'results/res_income', address, height=1000)
    else:
        results_path = screenshot_of_element(results_table, f'results/res_income_{index}', address, height=1000)

    driver.execute_script("window.scrollTo(0, 0)")
    return results_path


def results_res_rental(address, display_mode="for_sale"):
    escape()
    if display_mode == "for_sale":
        set_display(display_mode="display_for_sale")
    else:
        set_display(display_mode="display_marketing_to_realtors")
    sleep(1)

    escape()
    try:
        if display_mode == "for_sale":
            try:
                order_by_distance()
            except:
                print("WARNING: Can't Order by Distance. Maybe because there is not data")
        else:
            try:
                order_by_current_price()
                order_by_current_price()
            except:
                print("WARNING: Can't Order by Current Price. Maybe because there is not data")
    except Exception as e:
        print("WARNING: Cant order the data. \n")
    sleep(2)
    escape()

    sleep(2)
    # Screenshot
    results_table = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['results_table'])))
    results_path = screenshot_of_element(results_table, 'results/res_rental_' + display_mode, address, height=1000, width=1500)
    sleep(1)
    driver.execute_script("window.scrollTo(0, 0)")
    return results_path


def top_ten_links_on_google(address):
    driver.get(GOOGLE_URL)
    search_input = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, "//input[@type='text']")))
    search_input.send_keys(address)
    webdriver.ActionChains(driver).send_keys(Keys.ENTER).perform()

    links = driver.find_elements_by_xpath("//div[@class='yuRUbf']/a")[:10]
    links = [link.get_attribute("href") for link in links]
    return links


def extract_county_info(address):
    search_address_on_county(address)
    try:
        select_folio_number()
    except:
        print("INFO: Just one folio")
    return screenshots_of_county_info(address)


def search_address_on_county(address):
    driver.get(COUNTY_URL)

    search_input_id = "search_box"
    search_input = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.ID, search_input_id)))
    search_input.clear()
    search_input.send_keys(address)

    search_submit_id = "search_submit"
    search_submit = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.ID, search_submit_id)))
    search_submit.click()


def select_folio_number(index=1):
    WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, XPATHS['folio_number_link'])))
    folios_number = driver.find_elements_by_xpath(XPATHS['folio_number_link'])
    folio_number_link = folios_number[0]

    folio_number_link.click()


def screenshots_of_county_info(address):

    #print(driver.get_window_size())

    container = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, "//div[@class='container ng-scope']")))
    x_offset = (driver.get_window_size()['width'] - 1200)
    #print(x_offset)
    container_loc = container.location_once_scrolled_into_view
    subfolder = address.lower().replace(' ', '-').replace(',', '')

    sleep(1)
    # Property Info
    driver.execute_script("document.body.style.zoom='80%'")
    sleep(1)
    property_info_id = "property_info"
    property_info = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.ID, property_info_id)))
    sleep(1)

    location = property_info.location_once_scrolled_into_view
    size = property_info.size
    size['height'] = size['height'] * 0.8
    size['width'] = size['width'] * 0.8
    #location['x'] = x_offset + 20
    location['x'] = location['x'] * 0.8
    #location['y'] = location['y'] * 10

    sleep(1)
    property_info_path = screenshot_and_crop('county/' + subfolder, location, size, "property_info")
    sleep(1)
    driver.execute_script("document.body.style.zoom='100%'")
    driver.execute_script("window.scrollTo(0, 1000)")
    sleep(1)

    # Full Legal Desc
    full_legal_description_xpath = '//*[@id="contentScrollPoint"]/div[4]/div[4]/div[2]/div[2]/div/table'
    #"//table[@class='table table-condensed table-striped' and position()=3]"
    #"//div[@class='table-responsive ng-scope' and position()=4]"
    #"//div[@class='col-md-6' and position() = 2]/div[2]"

    full_legal_description = WebDriverWait(driver, 15).until(
        EC.element_to_be_clickable((By.XPATH, full_legal_description_xpath))
    )

    loc = full_legal_description.location_once_scrolled_into_view
    size = full_legal_description.size
    #loc['x'] = x_offset + 620
    #size['height'] = size['height'] * 1.20

    sleep(1)
    legal_info_path = screenshot_and_crop("county/" + subfolder, location=loc, size=size, filename="full_legal_info")
    #screenshot_of_element(full_legal_description, 'county/' + subfolder, "full_legal_info")

    # Taxable Info
    taxable_xpath = "//div[@class='col-md-6' and position() = 1]/div[2]"
    #"//div[@class='table-responsive ng-scope' and position()=2]"
    #"//div[@class='col-md-6' and position() = 1]/div[2]"
    taxable = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, taxable_xpath)))

    loc = taxable.location_once_scrolled_into_view
    size = taxable.size
    #loc['x'] = (x_offset/2) + 50
    #size['width'] = size['width'] * 1.25

    sleep(1)
    taxable_info_path = screenshot_and_crop("county/" + subfolder, location=loc,size=size, filename="taxable")
    #screenshot_of_element(taxable, 'county/' + subfolder, "taxable")

    # Sales Info
    sales_info_xpath = "//div[@class='row tabular_data' and not(@ng-show)]/div[@class='col-md-12']"
    #"//div[@class='table-responsive ng-scope' and position()=5]"
    #"//table[@class='table table-condensed table-striped' and position()=4]"
    # "//div[@class='row tabular_data' and not(@ng-show)]/div[@class='col-md-12']"
    sales_info = WebDriverWait(driver, 15).until(EC.element_to_be_clickable((By.XPATH, sales_info_xpath)))

    loc = sales_info.location_once_scrolled_into_view
    size = sales_info.size
    #loc['x'] = (x_offset/2) + 50
    #size['width'] = size['width'] * 1.2
    #size['height'] = size['height'] * 1.2

    sales_info_path = screenshot_and_crop("county/" + subfolder, location=loc, size=size, filename="sales_info")
    #screenshot_of_element(sales_info, 'county/' + subfolder, "sales_info")

    return {
        "property_info": property_info_path,
        "legal_info": legal_info_path,
        "taxable_info": taxable_info_path,
        "sales_info": sales_info_path
    }


def mls_extraction(address, baths, rooms, sqft_to):

    try:
        login()
    except:
        print("INFO: Login was not necessary")
    sleep(1)

    select_matrix_app()

    single_family_1 = single_family_search(address, baths, rooms, sqft_to)
    #printer.pprint(single_family_1)
    sleep(1)

    res_income_1 = res_income_search(address, baths, rooms, sqft_to)
    #printer.pprint(res_income_1)
    sleep(1)

    res_rental = res_rental_search(address, baths, rooms, sqft_to)
    #printer.pprint(res_rental)
    sleep(1)

    single_family_2 = single_family_search_2(address, baths, rooms, sqft_to)
    #printer.pprint(single_family_2)
    sleep(1)

    res_income_2 = res_income_search_2(address, baths, rooms, sqft_to)
    #printer.pprint(res_income_2)
    sleep(1)

    data = {
        "single_family_1": single_family_1,
        "res_income_1": res_income_1,
        "res_rental": res_rental,
        "single_family_2": single_family_2,
        "res_income_2": res_income_2
    }
    #print("Data:")
    #printer.pprint(data)
    return data

def extract(address, baths, rooms, sqft_to):
    # MLS Extraction
    criterias_results_paths = mls_extraction(address, baths, rooms, sqft_to)
    #printer.pprint(criterias_results_paths)

    # Extract the county info: property info, sales info, taxable info, legal info with screenshots
    try:
        county_info_paths = extract_county_info(address)
    except:
        county_info_paths = {
            "property_info": '',
            "legal_info": '',
            "taxable_info": '',
            "sales_info": ''
        }
    # Google the address, obtain the top 10 links
    top_google_links = top_ten_links_on_google(address)
    #printer.pprint(top_google_links)

    return top_google_links, criterias_results_paths, county_info_paths

def transform(address, mls_data, county_data):
    titles = [
        address.upper(),
        "PRIMEROS 10 ENLACES DE GOOGLE",
        "CRITERIO PARA COMPS FOR SALE SINGLE FAMILY",
        "COMPS FOR SALE SINGLE FAMILY",
        "CRITERIO PARA COMPS FOR SALE MULTI FAMILY",
        "COMPS FOR SALE MULTI FAMILY",
        "CRITERIO PARA COMPS FOR RENT",
        "COMPS FOR RENT (by Distance and Display For Sale)",
        "COMPS FOR RENT (by Higher Priced and Display MarketingToRealtor)",
        "COUNTY INFO: PROPERTY INFO",
        "COUNTY INFO: FULL LEGAL DESCRIPTION",
        "COUNTY INFO: TAXABLE",
        "COUNTY INFO: SALES INFO",
        "CRITERIO PARA COMPS FOR SALE SINGLE FAMILY (ACTIVE, PENDING Y ACTIVE WITH CONTRACT)",
        "COMPS FOR SALE SINGLE FAMILY (ACTIVE, PENDING Y ACTIVE WITH CONTRACT)",
        "CRITERIO PARA COMPS FOR SALE MULTI FAMILY (ACTIVE, PENDING Y ACTIVE WITH CONTRACT)",
        "COMPS FOR SALE MULTI FAMILY (ACTIVE, PENDING Y ACTIVE WITH CONTRACT)",
    ]
    slides_data = []
    slide_index = 2

    slides_data.append({'image_path': None, 'slide_index': 0, 'title': str(address).upper()})
    slides_data.append({'image_path': None, 'slide_index': 1, 'title': 'PRIMEROS 10 ENLACES DE GOOGLE'}, )
    county_keys = list(county_data.keys())
    mls_keys = list(mls_data.keys())

    similar_keys = mls_keys[-2:]
    mls_keys = mls_keys[:3]

    slide_index = 2
    for i in range(2, 16):
        if slide_index > 16:
            break

        if slide_index in [9, 10, 11, 12]:
            for county_key in county_keys:
                slide = {
                    "title": titles[slide_index],
                    "image_path": county_data[county_key],
                    "slide_index": slide_index
                }
                slides_data.append(slide)
                slide_index += 1
        else:
            keys_for_mls = similar_keys if slide_index > 9 else mls_keys

            for mls_key in keys_for_mls:
                for image_key in ['criteria', 'results']:
                    if mls_key == 'res_rental' and image_key == 'results':
                        for results_images in mls_data[mls_key][image_key]:
                            slide = {
                                "title": titles[slide_index],
                                "image_path": results_images,
                                "slide_index": slide_index
                            }
                            slides_data.append(slide)
                            slide_index += 1
                    else:
                        slide = {
                            "title": titles[slide_index],
                            "image_path": mls_data[mls_key][image_key],
                            "slide_index": slide_index
                        }
                        slides_data.append(slide)
                        slide_index += 1

                if slide_index == 9:
                    break

    #printer.pprint(slides_data)

    return slides_data

def load(address, google_links, slides_data):
    address_format = address.lower().replace(' ', '-').replace(',', '')

    prs = Presentation()
    i = 0
    for slide_data in slides_data:

        # Set layout
        if slide_data["slide_index"] == 0:
            layout = prs.slide_layouts[0]
        elif slide_data["slide_index"] == 1:
            layout = prs.slide_layouts[1]
        else:
            layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(layout)

        # Title for each slide
        title = slide.shapes.title
        title.text = slide_data['title']

        # Slide #1
        if slide_data["slide_index"] == 0:
            subtitle = slide.placeholders[1]
            subtitle.text = "MLS Matrix Bot"
        # Google Links Slide
        elif slide_data["slide_index"] == 1:
            body_shape = slide.shapes.placeholders[1]

            tf = body_shape.text_frame
            tf.text = 'Top 10:'
            # Google Links
            for link in google_links:
                p = tf.add_paragraph()
                p.text = link
                p.size = Pt(6)
                p.level = 1
                p.font.size = Pt(15)

        # Insert image if have a image path
        if slide_data["image_path"]:
            left = Inches(0.1)
            top = Inches(2) if len(slide_data['title']) > 40 else Inches(1.8)

            if slide_data["slide_index"] in [9, 10, 11]:
                height = Inches(5.5)
                width = None
                if slide_data['slide_index'] == 9:
                    top = Inches(1.5)
                    left = Inches(3)
                    width = Inches(3)
                    height = None
            else:
                width = Inches(9.9)
                height = None
            image_path = slide_data["image_path"]

            if os.path.isfile(image_path):
                pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)

        i += 1

    # Save
    prs.save(f'{address_format}.pptx')


def extract_transform_load(address, baths, rooms, sqft_to):

    # Extract
    google_links, mls_data, county_data = extract(address=address, baths=baths, rooms=rooms, sqft_to=sqft_to)
    # Transform
    slides_data = transform(address=address, mls_data=mls_data, county_data=county_data)
    # Load
    load(address=address, google_links=google_links, slides_data=slides_data)


if __name__ == "__main__":
    clear()

    # Create if not exists this folders
    verify_folder_exists('results')
    verify_folder_exists('criteria')
    verify_folder_exists('county')

    clear()
    while True:

        main_menu_action = prompt(main_menu, style=style)

        # Cerrar programa
        if main_menu_action['action'] == 'Cerrar':
            break
        elif main_menu_action['action'] == main_menu_actions[0]:

            address_info = prompt(address_info_input)
            address = address_info["address"]
            baths = address_info["baths"]
            rooms = address_info["rooms"]
            sqft_to = address_info["sqft_to"]

            baths = int(baths) if baths else baths
            rooms = int(rooms) if rooms else rooms
            sqft_to = int(sqft_to) if sqft_to else sqft_to

            driver = webdriver.Chrome(ChromeDriverManager().install(), options=options)
            driver.maximize_window()

            extract_transform_load(address=address, baths=baths, rooms=rooms, sqft_to=sqft_to)

        elif main_menu_action['action'] == main_menu_actions[1]:
            excel_info = prompt(excel_info_input)
            filename = excel_info['filename']
            wb = load_workbook(filename)
            sht = wb.active
            last_row = len(sht['B'])

            driver = webdriver.Chrome(ChromeDriverManager().install(), options=options)
            driver.maximize_window()

            for i in range(3, last_row):
                address = sht[STREET_COL + str(i)].value
                sqft_to = sht[SQFT_COL + str(i)].value
                baths = int(round(sht[BATHS_COL + str(i)].value))
                rooms = int(round(sht[BEDS_COL + str(i)].value))

                if address == None:
                    print()
                    print("*"*20)
                    print("\nMESSAGE: No Address at row", i, "\n")
                    print("*" * 20)
                    print()
                    break
                print(f"{i}.", "Address:", address, "SQFT:", sqft_to, "BATHS", baths, "ROOMS:",rooms)
                try:
                    extract_transform_load(address=address, sqft_to=sqft_to, baths=baths, rooms=rooms)
                except:
                    print("Error in Extracting Data")
                    traceback.print_exc()
                    break
